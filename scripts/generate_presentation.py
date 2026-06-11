from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import os

BASE_DIR = Path(__file__).resolve().parent.parent
REPORTS_DIR = BASE_DIR / 'reports'
CHARTS_DIR = REPORTS_DIR / 'eda_charts'

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    subtitle_box.text = subtitle

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(24)

def add_image_slide(prs, title, img_path):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if os.path.exists(img_path):
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(2), width=Inches(8))
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = f"[Image not found: {os.path.basename(img_path)}]"

def create_presentation():
    prs = Presentation()

    # 1. Title
    add_title_slide(prs, "Bluestock Mutual Fund Analytics", "Capstone Project Final Presentation\nBy Executive Data Team")

    # 2. Problem & Objective
    add_bullet_slide(prs, "Problem & Objective", [
        "Problem: Fragmented data across 40 schemes leading to poor visibility.",
        "Objective 1: Build a centralized Data Warehouse (Star Schema).",
        "Objective 2: Automate ETL pipelines to ingest live NAV data.",
        "Objective 3: Develop interactive dashboards for executive monitoring.",
        "Objective 4: Compute advanced performance metrics (Sharpe, Alpha, VaR)."
    ])

    # 3. Data Sources
    add_bullet_slide(prs, "Data Sources & Architecture", [
        "Internal Data: CSV dumps covering Demographics, Transactions, AUM, and Portfolio Holdings.",
        "External Data: AMFI Public API for live Daily NAV fetching.",
        "ETL Flow: Python/Pandas extraction -> Regex & Type casting -> SQLite Storage.",
        "Schema: Star Schema with Fact (NAV, Performance) and Dimension (Fund, Date) tables."
    ])

    # 4. Architecture Diagram
    # We will just use a bullet slide since we don't have a specific arch diagram image
    add_bullet_slide(prs, "Architecture Details", [
        "Extraction: Concurrent API calls to AMFI",
        "Transformation: Missing value imputation (ffill), Schema enforcement",
        "Loading: SQLite database for lightweight, portable analytical queries",
        "Reporting: Python generated PDFs, PPTX, and BI Dashboard integration"
    ])

    # 5. EDA Highlights 1
    add_image_slide(prs, "EDA Highlight: AUM Dominance", CHARTS_DIR / '2_aum_growth.png')

    # 6. EDA Highlights 2
    add_image_slide(prs, "EDA Highlight: Unstoppable SIP Inflows", CHARTS_DIR / '3_sip_inflows.png')

    # 7. Performance Metrics 1
    add_image_slide(prs, "Performance: Benchmark Comparison", BASE_DIR / 'benchmark_comparison_chart.png')

    # 8. Performance Metrics 2
    add_image_slide(prs, "Performance: 30-Day Rolling Sharpe", BASE_DIR / 'rolling_sharpe_chart.png')

    # 9. Dashboard Screenshot 1
    add_image_slide(prs, "Dashboard: Main Executive View", BASE_DIR / 'dashboard' / 'Page_1_Industry_Overview.png')

    # 10. Dashboard Screenshot 2
    add_image_slide(prs, "Dashboard: Performance Deep Dive", BASE_DIR / 'dashboard' / 'Page_2_Fund_Performance.png')

    # 11. Key Findings
    add_bullet_slide(prs, "Key Findings & Recommendations", [
        "Finding 1: 2023 Bull run significantly elevated equity fund baselines.",
        "Finding 2: B30 cities represent the highest growth vector.",
        "Rec 1: Target localized marketing campaigns in Tier-2/3 cities.",
        "Rec 2: Implement dynamic hedging for high-beta Small Cap funds."
    ])

    # 12. Thank You
    add_title_slide(prs, "Thank You", "Questions & Answers")

    # Output
    output_path = BASE_DIR / 'Bluestock_MF_Presentation.pptx'
    prs.save(str(output_path))
    print(f"Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    create_presentation()
